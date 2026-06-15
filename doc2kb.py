"""
doc2kb.py
---------
Convert enterprise documents with embedded text into clean Markdown
saved locally.

Stages:
    1. Extraction  - format-specific text extraction -> raw text
    2. Cleaning    - removes noise (page numbers, headers, footers,
                     broken spacing, URLs, bibliography)
    3. Metadata    - prepends YAML front-matter
    4. Save        - writes <filename>_clean.md next to the input file

INSTALL:
    pip install pymupdf python-docx python-pptx beautifulsoup4

USAGE:
    python doc2kb.py report.pdf
    python doc2kb.py *.pdf *.docx
    python doc2kb.py report.pdf --doc_type supply_chain_ops
    python doc2kb.py report.pdf --dry_run
"""

import argparse
import logging
import os
import re
import unicodedata
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path

logging.basicConfig(level=logging.INFO, format="%(levelname)s | %(message)s")
log = logging.getLogger(__name__)


# -----------------------------------------------------------------------------
# STAGE 1: EXTRACTION
# -----------------------------------------------------------------------------

def extract_pdf(path: str) -> str:
    try:
        import fitz
    except ImportError as exc:
        raise ImportError("PyMuPDF not installed. Run: pip install pymupdf") from exc

    doc = fitz.open(path)
    try:
        pages = [page.get_text() for page in doc]
    finally:
        doc.close()

    return "\n\n".join(pages)


def extract_docx(path: str) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise ImportError(
            "python-docx not installed. Run: pip install python-docx"
        ) from exc

    doc = Document(path)
    return "\n".join(
        paragraph.text
        for paragraph in doc.paragraphs
        if paragraph.text.strip()
    )


def extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ImportError(
            "python-pptx not installed. Run: pip install python-pptx"
        ) from exc

    prs = Presentation(path)
    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text.append(shape.text)

    return "\n\n".join(text)


def extract_html(path: str) -> str:
    try:
        from bs4 import BeautifulSoup
    except ImportError as exc:
        raise ImportError(
            "beautifulsoup4 not installed. Run: pip install beautifulsoup4"
        ) from exc

    with open(path, encoding="utf-8") as f:
        soup = BeautifulSoup(f.read(), "html.parser")

    return soup.get_text("\n")


def extract_to_markdown(input_path: str) -> str:
    """
    Extract plain text from supported document types.
    """
    path = Path(input_path)
    suffix = path.suffix.lower()
    extractor_map = {
        ".pdf": ("PyMuPDF", extract_pdf),
        ".docx": ("python-docx", extract_docx),
        ".pptx": ("python-pptx", extract_pptx),
        ".html": ("BeautifulSoup", extract_html),
        ".htm": ("BeautifulSoup", extract_html),
    }

    if suffix not in extractor_map:
        supported = ", ".join(sorted(extractor_map))
        raise ValueError(
            f"Unsupported file type '{suffix or '[no extension]'}'. "
            f"Supported types: {supported}"
        )

    extractor_name, extractor = extractor_map[suffix]
    log.info(f"[1/3] Extracting: {path.name}")
    log.info(f"Extractor: {extractor_name}")
    text = extractor(str(path))
    log.info(f"Characters extracted: {len(text):,}")
    return text


# -----------------------------------------------------------------------------
# STAGE 2: CLEANING
# -----------------------------------------------------------------------------

_REPEAT_THRESHOLD = 0.15

_STOP_SECTIONS = re.compile(
    r"^#{1,4}\s*(References|Bibliography|Works\s+Cited|Further\s+Reading"
    r"|Notes|Endnotes|Acknowledgements?|Appendix|Index)\s*$",
    re.IGNORECASE | re.MULTILINE,
)
_URL_RE = re.compile(
    r"https?://\S+|ftp://\S+|www\.\S+|doi\.org/\S+|10\.\d{4,}/\S+",
    re.IGNORECASE,
)
_PAGE_NUM_RE = re.compile(
    r"^\s*[-–]?\s*(Page\s*)?\d{1,4}\s*[-–]?\s*$",
    re.IGNORECASE | re.MULTILINE,
)
_BROKEN_SPACE_RE = re.compile(r"\b([a-zA-Z]) ([a-zA-Z]{2,})\b")


def clean_markdown(md: str) -> str:
    """
    Stage 1  - Unicode normalization
    Stage 2  - Ligature + typographic character repair
    Stage 3  - Broken spacing fix  (s upply -> supply)
    Stage 4  - URL removal
    Stage 5  - Page number removal
    Stage 6  - Repeated header/footer removal
    Stage 7  - Bibliography/references section cut
    Stage 8  - Hyphenated line-break repair
    Stage 9  - Whitespace normalization
    Stage 10 - Empty heading removal
    """
    log.info("[2/3] Cleaning")

    # Stage 1
    md = unicodedata.normalize("NFC", md)

    # Stage 2
    ligatures = {
        "\ufb00": "ff",
        "\ufb01": "fi",
        "\ufb02": "fl",
        "\ufb03": "ffi",
        "\ufb04": "ffl",
        "\ufb06": "st",
        "\u2019": "'",
        "\u2018": "'",
        "\u201c": '"',
        "\u201d": '"',
        "\u2013": "-",
        "\u2014": "--",
    }
    for src, dst in ligatures.items():
        md = md.replace(src, dst)

    # Stage 3 - multiple passes until stable
    for _ in range(8):
        fixed = _BROKEN_SPACE_RE.sub(lambda m: m.group(1) + m.group(2), md)
        if fixed == md:
            break
        md = fixed

    # Stage 4
    md = _URL_RE.sub("", md)

    # Stage 5
    md = _PAGE_NUM_RE.sub("", md)

    # Stage 6
    md = _remove_repeated_lines(md)

    # Stage 7
    stop = _STOP_SECTIONS.search(md)
    if stop:
        md = md[:stop.start()].rstrip()
        log.info("      Removed references/bibliography section")

    # Stage 8
    md = re.sub(r"-\s*\n\s*([a-z])", r"\1", md)

    # Stage 9
    md = re.sub(r"\n{3,}", "\n\n", md)
    md = "\n".join(line.rstrip() for line in md.splitlines())

    # Stage 10
    md = re.sub(r"^#{1,6}\s*$", "", md, flags=re.MULTILINE)

    md = md.strip()
    log.info(f"      After cleaning: {len(md.split()):,} words")
    return md


def _remove_repeated_lines(md: str) -> str:
    lines = md.splitlines()
    if len(lines) < 10:
        return md

    freq: Counter = Counter()
    for line in lines:
        stripped = line.strip()
        if not stripped or stripped.startswith("#") or len(stripped) < 8:
            continue
        normalized = re.sub(r"[*_`]", "", stripped).lower()
        freq[normalized] += 1

    threshold = max(3, int(len(lines) * _REPEAT_THRESHOLD))
    repeated = {text for text, count in freq.items() if count >= threshold}

    if repeated:
        log.info(f"      Removed {len(repeated)} repeated header/footer line(s)")

    return "\n".join(
        line
        for line in lines
        if re.sub(r"[*_`]", "", line.strip()).lower() not in repeated
    )


# -----------------------------------------------------------------------------
# STAGE 3: METADATA
# -----------------------------------------------------------------------------

def add_metadata(md: str, source_file: str, doc_type: str) -> str:
    """Prepend YAML front-matter to the clean Markdown."""
    title = Path(source_file).stem.replace("_", " ").replace("-", " ").title()
    ingested_at = datetime.now(timezone.utc).strftime("%Y-%m-%d")

    front_matter = (
        f"---\n"
        f"title: {title}\n"
        f"source_file: {source_file}\n"
        f"doc_type: {doc_type}\n"
        f"ingested_at: {ingested_at}\n"
        f"---\n\n"
    )
    return front_matter + md


# -----------------------------------------------------------------------------
# STAGE 4: SAVE LOCALLY
# -----------------------------------------------------------------------------

def save_locally(md: str, input_path: str) -> str:
    """
    Save <filename>_clean.md in the same directory as the input file.
    Returns the output path.
    """
    input_path = Path(input_path).resolve()
    out_path = input_path.parent / f"{input_path.stem}_clean.md"
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(md)
    log.info(f"[3/3] Saved -> {out_path}")
    return str(out_path)


# -----------------------------------------------------------------------------
# ORCHESTRATOR
# -----------------------------------------------------------------------------

def process_document(
    input_path: str,
    doc_type: str,
    dry_run: bool = False,
) -> str | None:
    source_file = Path(input_path).name

    log.info(f"\n{'-' * 50}")
    log.info(f"  {source_file}")
    log.info(f"{'-' * 50}")

    try:
        raw_md = extract_to_markdown(input_path)
    except Exception as e:
        log.error(f"Extraction failed: {e}")
        return None

    clean_md = clean_markdown(raw_md)
    final_md = add_metadata(clean_md, source_file=source_file, doc_type=doc_type)

    if dry_run:
        sep = "-" * 50
        print(f"\n{sep}\nDRY RUN - {source_file}\n{sep}")
        print(final_md[:2500])
        if len(final_md) > 2500:
            print(f"\n... ({len(final_md) - 2500:,} more characters)")
        print(f"{sep}\n")
        return None

    return save_locally(final_md, input_path)


# -----------------------------------------------------------------------------
# CLI
# -----------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Convert PDF/DOCX/PPTX/HTML -> clean Markdown, saved locally.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python doc2kb.py report.pdf
  python doc2kb.py *.pdf *.docx
  python doc2kb.py report.pdf --doc_type supply_chain_ops
  python doc2kb.py report.pdf --dry_run
        """,
    )
    parser.add_argument("files", nargs="+", help="PDF, DOCX, PPTX, or HTML file(s)")
    parser.add_argument(
        "--doc_type",
        default="business_knowledge",
        help="Document category for metadata (default: business_knowledge)\n"
        "e.g. supply_chain_ops, demand_planning, logistics, hr_policy",
    )
    parser.add_argument(
        "--dry_run",
        action="store_true",
        help="Preview cleaned output without saving",
    )
    parser.add_argument("--verbose", "-v", action="store_true")
    args = parser.parse_args()

    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)

    results = []
    for f in args.files:
        if not os.path.isfile(f):
            log.warning(f"Skipping (not found): {f}")
            continue
        out = process_document(f, doc_type=args.doc_type, dry_run=args.dry_run)
        if out:
            results.append(out)

    if not args.dry_run:
        print(f"\nProcessed {len(results)}/{len(args.files)} file(s)")
        for r in results:
            print(f"  -> {r}")


if __name__ == "__main__":
    main()
